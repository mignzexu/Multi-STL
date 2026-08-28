# pyright: basic
import json
import struct
import zipfile
from pathlib import Path

import imageio
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
REPORTS = ROOT / "reports"
ASSETS = REPORTS / "assets"
OUT_PPTX = REPORTS / "multi_stl_gvbf_experiment_report_editable.pptx"
MANIFEST = REPORTS / "multi_stl_gvbf_experiment_manifest.json"
FONT_PATH = "/tmp/opencode/fonts/LXGWWenKai-Regular.ttf"

COLORS = {
    "ink": RGBColor(0x0F, 0x17, 0x2A),
    "ink2": RGBColor(0x1E, 0x29, 0x3B),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "accent": RGBColor(0xF9, 0x73, 0x16),
    "muted": RGBColor(0x64, 0x74, 0x8B),
}

EXP_ORDER = [
    ("biflow", "biflow_w3s", "BiFlow 模式", "联合前向/后向轨迹与条件噪声建模", "7C3AED"),
    ("flow", "flow_w3s", "Flow 模式", "确定性流轨迹用于风速序列预报", "0EA5E9"),
    ("condiff", "condiff_w3s", "ConDiff 模式", "基于条件扩散的 GVBF 动力学校正", "10B981"),
]
THRESHOLDS = [8.0, 13.9, 20.8]
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def hex_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def gif_to_avi(gif_path, avi_path, fps=10):
    reader = imageio.get_reader(str(gif_path))
    frames = [reader.get_data(i) for i in range(len(reader))]
    reader.close()
    w, h = frames[0].shape[1], frames[0].shape[0]
    n = len(frames)
    duration = int(1000000 / fps)
    fd_list = []
    max_fd = 0
    for frame in frames:
        img = Image.fromarray(frame).convert("RGB")
        buf = __import__("io").BytesIO()
        img.save(buf, format="JPEG", quality=85)
        fd = buf.getvalue()
        fd_list.append(fd)
        max_fd = max(max_fd, len(fd))
    avi = bytearray()
    avi += b"RIFF\x00\x00\x00\x00AVI "
    hdrl = bytearray()
    hdrl += b"LIST\x00\x00\x00\x00hdrl"
    strh = struct.pack("<IIIIIIIIIIHHHH", 0, 0, 0, 1, fps, 0, n, max_fd, w, h, 0, 0, 0, 0)
    strf = struct.pack("<IiiHHIIiiII", 40, w, h, 1, 24, 0, w * h * 3, 0, 0, 0, 0)
    strl = bytearray(b"LIST\x00\x00\x00\x00strl")
    strl += b"strh" + struct.pack("<I", len(strh)) + strh
    strl += b"strf" + struct.pack("<I", len(strf)) + strf
    if len(strl) % 2: strl += b"\x00"
    avih = struct.pack("<IIIIIIIIIIIIII", duration, max_fd * fps, 0, 0x10, n, max_fd, w, h, 0, 0, 0, 0, 0, 0)
    hdrl += b"avih" + struct.pack("<I", len(avih)) + avih
    hdrl += strl
    if len(hdrl) % 2: hdrl += b"\x00"
    hdrl[4:8] = struct.pack("<I", len(hdrl) - 8)
    avi += hdrl
    movi = bytearray(b"LIST\x00\x00\x00\x00movi")
    for fd in fd_list:
        padded = fd + b"\x00" * (len(fd) % 2)
        movi += b"00dc" + struct.pack("<I", len(fd)) + padded
    if len(movi) % 2: movi += b"\x00"
    movi[4:8] = struct.pack("<I", len(movi) - 8)
    avi += movi
    avi[4:8] = struct.pack("<I", len(avi) - 8)
    avi_path.write_bytes(bytes(avi))


def add_rect(s, l, t, w, h, fill=None, lc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    if lc: sh.line.color.rgb = lc
    else: sh.line.fill.background()
    return sh


def add_rrect(s, l, t, w, h, fill=None, lc=None, lw=Pt(1)):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    if lc: sh.line.color.rgb = lc; sh.line.width = lw
    else: sh.line.fill.background()
    return sh


def tb(s, l, t, w, h, txt, sz=Pt(18), clr=COLORS["ink"], bold=False, al=PP_ALIGN.LEFT):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.alignment = al; p.space_before = Pt(0); p.space_after = Pt(0)
    r = p.add_run(); r.text = txt; r.font.size = sz; r.font.color.rgb = clr
    r.font.bold = bold; r.font.name = "LXGW WenKai"
    return bx


def gen_chart(exp, path, W=900, H=360):
    import math
    ft = ImageFont.truetype(FONT_PATH, 18)
    ft_s = ImageFont.truetype(FONT_PATH, 14)
    ft_t = ImageFont.truetype(FONT_PATH, 20)
    img = Image.new("RGB", (W, H), "white")
    d = ImageDraw.Draw(img)
    tr = [(p["epoch"], p["train"]) for p in exp["epochs"] if p["train"] is not None]
    va = [(p["epoch"], p["valid"]) for p in exp["epochs"] if p["valid"] is not None]
    vals = [v for _, v in tr + va]
    if not vals:
        d.text((24, 120), "训练历史不可用", font=ft, fill="gray")
        img.save(path, dpi=(144, 144)); return
    ym, yx = min(vals), max(vals)
    if math.isclose(ym, yx): ym *= 0.95; yx *= 1.05
    pad = (yx - ym) * 0.08; ym -= pad; yx += pad
    x0, y0, x1, y1 = 58, 28, W - 28, H - 48
    d.rectangle((x0, y0, x1, y1), outline="#CBD5E1", width=2)
    for i in range(1, 4):
        yy = y0 + (y1 - y0) * i / 4; d.line((x0, yy, x1, yy), fill="#E2E8F0", width=1)
    me = max(p["epoch"] for p in exp["epochs"])
    def xy(pt):
        e, v = pt
        return (x0 + (x1 - x0) * (e - 1) / max(1, me - 1), y1 - (y1 - y0) * (v - ym) / (yx - ym))
    for ser, col, nm, yl in [(tr, "#F97316", "训练", 14), (va, "#" + exp["color"], "验证", 42)]:
        if len(ser) > 1: d.line([xy(p) for p in ser], fill=col, width=4)
        d.rounded_rectangle((W - 150, yl, W - 130, yl + 20), 4, fill=col)
        d.text((W - 120, yl - 2), nm, font=ft_s, fill="#0F172A")
    d.text((24, 8), "损失曲线", font=ft_t, fill="#0F172A")
    d.text((x0, H - 35), "epoch 1", font=ft_s, fill="#64748B")
    d.text((x1 - 82, H - 35), f"epoch {me}", font=ft_s, fill="#64748B")
    img.save(path, dpi=(144, 144))


def load_exp(slug, folder, label, tag, color):
    root = ROOT / "work_dirs" / folder
    cfg = json.loads((root / "obj_config.json").read_text())
    res = json.loads((root / "result.json").read_text())["w10"]
    met = json.loads((root / "metrics.json").read_text())
    eps = []
    for k, v in met.items():
        if k.isdigit():
            eps.append({"epoch": int(k), "train": v.get("train", {}).get("loss"), "valid": v.get("valid", {}).get("loss")})
    eps.sort(key=lambda x: x["epoch"])
    vp = [p for p in eps if p["valid"] is not None]
    best = min(vp, key=lambda x: x["valid"])
    ckpts = sorted((root / "model").glob("epoch*.ckpt"))
    gif = root / "vis" / "2017_0_browse.gif"
    if not gif.exists():
        cs = sorted((root / "vis").glob("*_browse.gif"))
        gif = cs[0] if cs else None
    return {
        "slug": slug, "folder": folder, "label": label, "tag": tag, "color": color,
        "cfg": {k: cfg.get(k) for k in ["gvbf_mode","gvbf_model_config","dataset","data_config","total_seq","img_size","in_category","out_category","learning_rate","batch_size","seed"]},
        "result": res, "epochs": eps, "best_epoch": best["epoch"], "best_val": best["valid"],
        "final_val": vp[-1]["valid"] if vp else None, "ckpt": ckpts[0].name if ckpts else None,
        "gif": str(gif) if gif else None,
    }


def load_all():
    items = [load_exp(*s) for s in EXP_ORDER]
    for m in ["mae", "rmse", "mse"]:
        for i, it in enumerate(sorted(items, key=lambda x: x["result"][m]), 1):
            it.setdefault("rank", {})[m] = i
    for i, it in enumerate(sorted(items, key=lambda x: sum(x["result"]["hss"]) / 3, reverse=True), 1):
        it.setdefault("rank", {})["hss_avg"] = i
    MANIFEST.write_text(json.dumps(items, indent=2, ensure_ascii=False))
    return items


def build_setup(prs, exp):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    ac = hex_rgb(exp["color"])
    c = exp["cfg"]
    add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(1.25), fill=COLORS["ink"])
    add_rect(sl, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=ac)
    tb(sl, Inches(0.6), Inches(0.15), Inches(10), Inches(0.65), f"{exp['label']}: 实验设置", Pt(36), COLORS["white"], bold=True)
    tb(sl, Inches(0.6), Inches(0.75), Inches(10), Inches(0.4), exp["tag"], Pt(18), hex_rgb("CBD5E1"))
    add_rrect(sl, Inches(0.6), Inches(1.6), Inches(5.6), Inches(5.3), fill=COLORS["white"], lc=hex_rgb("E2E8F0"))
    tb(sl, Inches(0.9), Inches(1.8), Inches(5), Inches(0.5), "实验配置", Pt(28), COLORS["ink"], bold=True)
    rows = [("模式", c.get("gvbf_mode")), ("GVBF 模型配置", c.get("gvbf_model_config")),
            ("数据集", f"{c.get('dataset')} / {c.get('data_config')}"),
            ("变量", f"{c.get('in_category')} -> {c.get('out_category')}"),
            ("序列", f"{c.get('total_seq')[0]} 输入 + {c.get('total_seq')[1]} 预测帧"),
            ("网格尺寸", f"{c.get('img_size')[0]} x {c.get('img_size')[1]}"),
            ("优化器", f"lr={c.get('learning_rate')}, batch={c.get('batch_size')}, seed={c.get('seed')}"),
            ("最优检查点", exp.get("ckpt"))]
    y = Inches(2.5)
    for k, v in rows:
        tb(sl, Inches(0.95), y, Inches(1.8), Inches(0.4), k, Pt(14), COLORS["muted"], bold=True)
        tb(sl, Inches(2.6), y, Inches(3.3), Inches(0.4), str(v), Pt(16), COLORS["ink"])
        y += Inches(0.5)
    add_rrect(sl, Inches(6.7), Inches(1.6), Inches(6.1), Inches(2.8), fill=COLORS["ink2"])
    tb(sl, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.5), "本实验测试内容", Pt(28), COLORS["white"], bold=True)
    by = Inches(2.5)
    for b in ["三种模式均在 WeatherBench wind_3s 任务上进行。", "基于 12 帧历史风速数据预报未来 12 帧。",
              "评估指标：MAE/MSE/RMSE 及阈值指标（8.0, 13.9, 20.8）。"]:
        tb(sl, Inches(7.3), by, Inches(5.2), Inches(0.5), b, Pt(16), hex_rgb("E2E8F0"))
        by += Inches(0.5)
    add_rrect(sl, Inches(6.7), Inches(4.7), Inches(6.1), Inches(2.3), fill=COLORS["white"], lc=hex_rgb("E2E8F0"))
    tb(sl, Inches(7.0), Inches(4.85), Inches(5.5), Inches(0.4), "代表性预测动画：2017_0", Pt(18), COLORS["ink2"], bold=True)
    avi_path = ASSETS / f"{exp['slug']}_pred.avi"
    if exp["gif"]:
        gif_to_avi(exp["gif"], avi_path)
        sl.shapes.add_movie(str(avi_path), Inches(7.0), Inches(5.3), Inches(5.5), Inches(1.5))
    tb(sl, Inches(0.6), Inches(7.1), Inches(12), Inches(0.3), f"数据来源：work_dirs/{exp['folder']}", Pt(11), COLORS["ink2"])


def build_results(prs, exp):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    ac = hex_rgb(exp["color"])
    r = exp["result"]
    add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(1.25), fill=COLORS["ink"])
    add_rect(sl, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=ac)
    tb(sl, Inches(0.6), Inches(0.15), Inches(10), Inches(0.65), f"{exp['label']}: 实验结果", Pt(36), COLORS["white"], bold=True)
    tb(sl, Inches(0.6), Inches(0.75), Inches(10), Inches(0.4), "测试指标与训练动态", Pt(18), hex_rgb("CBD5E1"))
    cards = [("MAE", f"{r['mae']:.3f}", Inches(0.6), exp["rank"].get("mae")),
             ("RMSE", f"{r['rmse']:.3f}", Inches(2.8), exp["rank"].get("rmse")),
             ("最优验证损失", f"{exp['best_val']:.4f}", Inches(5.0), None),
             ("最优 Epoch", str(exp["best_epoch"]), Inches(7.2), None)]
    for label, val, x, rk in cards:
        add_rrect(sl, x, Inches(1.6), Inches(2.0), Inches(1.15), fill=COLORS["white"], lc=ac, lw=Pt(2))
        tb(sl, x + Inches(0.2), Inches(1.7), Inches(1.6), Inches(0.3), label, Pt(12), ac, bold=True)
        tb(sl, x + Inches(0.2), Inches(2.05), Inches(1.6), Inches(0.5), val, Pt(28), COLORS["ink"], bold=True)
        if rk: tb(sl, x + Inches(0.2), Inches(2.45), Inches(1.6), Inches(0.2), f"排名 {rk}/3（越低越好）", Pt(9), COLORS["muted"])
    add_rrect(sl, Inches(9.4), Inches(1.6), Inches(3.4), Inches(1.15), fill=COLORS["ink2"])
    avg_hss = sum(r["hss"]) / 3
    tb(sl, Inches(9.6), Inches(1.7), Inches(3), Inches(0.3), "平均 HSS", Pt(14), hex_rgb("CBD5E1"), bold=True)
    tb(sl, Inches(9.6), Inches(2.0), Inches(3), Inches(0.5), f"{avg_hss:.3f}", Pt(28), COLORS["white"], bold=True)
    tb(sl, Inches(11.2), Inches(2.45), Inches(1.5), Inches(0.2), f"排名 {exp['rank']['hss_avg']}/3", Pt(12), hex_rgb("E2E8F0"), bold=True)
    add_rrect(sl, Inches(0.6), Inches(3.1), Inches(4.8), Inches(4.0), fill=COLORS["white"], lc=hex_rgb("E2E8F0"))
    tb(sl, Inches(0.9), Inches(3.25), Inches(4), Inches(0.4), "阈值指标", Pt(20), COLORS["ink"], bold=True)
    tb(sl, Inches(0.9), Inches(3.65), Inches(4), Inches(0.3), "风速阈值下的 CSI / POD / FAR / HSS", Pt(12), COLORS["muted"])
    cx = [Inches(0.95), Inches(1.85), Inches(2.75), Inches(3.65), Inches(4.35)]
    y = Inches(4.2)
    for xx, h in zip(cx, ["阈值", "CSI", "POD", "FAR", "HSS"]):
        tb(sl, xx, y, Inches(0.85), Inches(0.3), h, Pt(13), ac, bold=True)
    y += Inches(0.45)
    for i, thr in enumerate(THRESHOLDS):
        for xx, v in zip(cx, [f"{thr:g}", f"{r['csi'][i]:.3f}", f"{r['pod'][i]:.3f}", f"{r['far'][i]:.3f}", f"{r['hss'][i]:.3f}"]):
            tb(sl, xx, y, Inches(0.85), Inches(0.35), v, Pt(14), COLORS["ink"])
        y += Inches(0.45)
    tb(sl, Inches(0.9), Inches(6.4), Inches(4), Inches(0.3), "观察结论", Pt(18), COLORS["ink"], bold=True)
    note = {"flow": "三个实验中综合最优：MAE/RMSE 最低，平均 HSS 最高。",
            "condiff": "检测指标接近 Flow，但误差幅度仍然较高。"}.get(exp["slug"], "验证集最优出现较早（epoch 4），最终测试指标弱于其他模式。")
    tb(sl, Inches(0.9), Inches(6.75), Inches(4.2), Inches(0.5), note, Pt(14), COLORS["ink"])
    add_rrect(sl, Inches(5.8), Inches(3.1), Inches(7.0), Inches(4.0), fill=COLORS["white"], lc=hex_rgb("E2E8F0"))
    tb(sl, Inches(6.1), Inches(3.25), Inches(6), Inches(0.4), "训练动态", Pt(20), COLORS["ink"], bold=True)
    cp = ASSETS / f"{exp['slug']}_loss_chart.png"
    gen_chart(exp, cp)
    sl.shapes.add_picture(str(cp), Inches(6.1), Inches(3.8), Inches(5.8), Inches(2.3))
    tb(sl, Inches(6.1), Inches(6.2), Inches(6.5), Inches(0.4), f"最终验证损失：{exp['final_val']:.4f} | 检查点：{exp['ckpt']}", Pt(14), COLORS["ink2"])
    tb(sl, Inches(0.6), Inches(7.1), Inches(12), Inches(0.3), f"指标来源：work_dirs/{exp['folder']}/result.json 与 metrics.json", Pt(11), COLORS["ink2"])


def build_future(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(1.25), fill=COLORS["ink"])
    add_rect(sl, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=COLORS["accent"])
    tb(sl, Inches(0.6), Inches(0.15), Inches(10), Inches(0.65), "未来工作", Pt(36), COLORS["white"], bold=True)
    tb(sl, Inches(0.6), Inches(0.75), Inches(10), Inches(0.4), "后续研究方向与计划", Pt(18), hex_rgb("CBD5E1"))
    items = [
        ("01", "本地数据集验证", "将模型在本地气象站点数据集上进行验证，评估模型在真实观测数据上的泛化能力和预报精度。"),
        ("02", "25 雷达数据处理", "完成 25 雷达数据的预处理流水线，包括质量控制、时空对齐和标准化，为多源融合实验提供数据基础。"),
        ("03", "WeatherBench 新 Baseline", "在 WeatherBench 基准中继续验证新的 Baseline 模型，对比 GVBF 各模式与主流方法的性能差异。"),
    ]
    y = Inches(1.8)
    for num, title, desc in items:
        add_rrect(sl, Inches(0.8), y, Inches(11.7), Inches(1.5), fill=COLORS["white"], lc=COLORS["accent"], lw=Pt(2))
        add_rrect(sl, Inches(1.1), y + Inches(0.25), Inches(0.7), Inches(0.7), fill=COLORS["accent"])
        tb(sl, Inches(1.1), y + Inches(0.3), Inches(0.7), Inches(0.6), num, Pt(28), COLORS["white"], bold=True, al=PP_ALIGN.CENTER)
        tb(sl, Inches(2.1), y + Inches(0.2), Inches(10), Inches(0.45), title, Pt(24), COLORS["ink"], bold=True)
        tb(sl, Inches(2.1), y + Inches(0.7), Inches(10), Inches(0.6), desc, Pt(16), COLORS["muted"])
        y += Inches(1.8)
    tb(sl, Inches(0.6), Inches(7.1), Inches(12), Inches(0.3), "Multi-STL 降水临近预报研究计划", Pt(11), COLORS["ink2"])


def main():
    ASSETS.mkdir(parents=True, exist_ok=True)
    exps = load_all()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for exp in exps:
        build_setup(prs, exp)
        build_results(prs, exp)
    build_future(prs)
    prs.save(str(OUT_PPTX))
    with zipfile.ZipFile(OUT_PPTX) as zf:
        assert zf.testzip() is None
        assert len([n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]) == 7
    data = json.loads(MANIFEST.read_text())
    by_slug = {i["slug"]: i for i in data}
    for slug, folder in [("biflow","biflow_w3s"),("flow","flow_w3s"),("condiff","condiff_w3s")]:
        src = json.loads((ROOT/"work_dirs"/folder/"result.json").read_text())["w10"]
        for m in ["mae","mse","rmse"]:
            assert by_slug[slug]["result"][m] == src[m]
    print(OUT_PPTX)

if __name__ == "__main__":
    main()
